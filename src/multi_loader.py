"""
Multi-format document loader supporting PDF, DOCX, PPTX, TXT, MD, CSV, HTML, URLs.
Includes OCR support for scanned PDFs and images.
"""

import hashlib
import re
import logging
from pathlib import Path
from typing import Dict, List, Optional
from dataclasses import dataclass

logger = logging.getLogger("rag.dataloader")


@dataclass
class LoadedDocument:
    doc_id: str
    text: str
    source: str
    doc_type: str
    metadata: dict = None

    def __post_init__(self):
        if self.metadata is None:
            self.metadata = {}


def _hash(text: str) -> str:
    return hashlib.md5(text.encode()).hexdigest()[:12]


def load_txt(path: str) -> List[LoadedDocument]:
    text = Path(path).read_text(encoding="utf-8", errors="replace")
    return [LoadedDocument(
        doc_id=_hash(text),
        text=text,
        source=path,
        doc_type="txt",
    )]


def load_markdown(path: str) -> List[LoadedDocument]:
    text = Path(path).read_text(encoding="utf-8", errors="replace")
    return [LoadedDocument(
        doc_id=_hash(text),
        text=text,
        source=path,
        doc_type="markdown",
    )]


def load_csv(path: str) -> List[LoadedDocument]:
    import pandas as pd
    df = pd.read_csv(path)
    docs = []
    for i, row in df.iterrows():
        text = " | ".join(f"{col}: {val}" for col, val in row.items() if str(val) != "nan")
        docs.append(LoadedDocument(
            doc_id=f"{_hash(path)}_row_{i}",
            text=text,
            source=path,
            doc_type="csv_row",
            metadata={"row_index": i, "columns": list(df.columns)},
        ))
    return docs


def load_html(path: str) -> List[LoadedDocument]:
    from html.parser import HTMLParser

    class TextExtractor(HTMLParser):
        def __init__(self):
            super().__init__()
            self.text = []
            self._skip = False

        def handle_starttag(self, tag, attrs):
            if tag in ('script', 'style'):
                self._skip = True

        def handle_endtag(self, tag):
            if tag in ('script', 'style'):
                self._skip = False
            elif tag in ('p', 'div', 'br', 'li', 'h1', 'h2', 'h3', 'h4'):
                self.text.append('\n')

        def handle_data(self, data):
            if not self._skip:
                self.text.append(data)

    html = Path(path).read_text(encoding="utf-8", errors="replace")
    extractor = TextExtractor()
    extractor.feed(html)
    text = " ".join(extractor.text)
    text = re.sub(r'\s+', ' ', text).strip()
    return [LoadedDocument(
        doc_id=_hash(text),
        text=text,
        source=path,
        doc_type="html",
    )]


def load_docx(path: str) -> List[LoadedDocument]:
    try:
        import docx
        doc = docx.Document(path)
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return [LoadedDocument(
            doc_id=_hash(text),
            text=text,
            source=path,
            doc_type="docx",
        )]
    except ImportError:
        logger.warning("python-docx not installed, cannot load DOCX")
        return []


def load_pptx(path: str) -> List[LoadedDocument]:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        texts = []
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                texts.append(f"[Slide {slide_num + 1}] " + " ".join(slide_text))
        text = "\n\n".join(texts)
        return [LoadedDocument(
            doc_id=_hash(text),
            text=text,
            source=path,
            doc_type="pptx",
        )]
    except ImportError:
        logger.warning("python-pptx not installed, cannot load PPTX")
        return []


def load_url(url: str) -> List[LoadedDocument]:
    try:
        import urllib.request
        from html.parser import HTMLParser

        class TextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.text = []
                self._skip = False
            def handle_starttag(self, tag, attrs):
                if tag in ('script', 'style', 'nav', 'footer', 'header'):
                    self._skip = True
            def handle_endtag(self, tag):
                if tag in ('script', 'style', 'nav', 'footer', 'header'):
                    self._skip = False
                elif tag in ('p', 'div', 'br', 'li', 'h1', 'h2', 'h3', 'h4'):
                    self.text.append('\n')
            def handle_data(self, data):
                if not self._skip:
                    self.text.append(data)

        req = urllib.request.Request(url, headers={'User-Agent': 'CorrectiveRAG/1.0'})
        with urllib.request.urlopen(req, timeout=30) as resp:
            html = resp.read().decode('utf-8', errors='replace')

        extractor = TextExtractor()
        extractor.feed(html)
        text = " ".join(extractor.text)
        text = re.sub(r'\s+', ' ', text).strip()
        return [LoadedDocument(
            doc_id=_hash(url + text),
            text=text,
            source=url,
            doc_type="url",
            metadata={"url": url},
        )]
    except Exception as e:
        logger.error(f"Failed to load URL {url}: {e}")
        return []


def load_image_ocr(path: str) -> List[LoadedDocument]:
    """OCR support for images using pytesseract."""
    try:
        from PIL import Image
        import pytesseract
        img = Image.open(path)
        text = pytesseract.image_to_string(img)
        if text.strip():
            return [LoadedDocument(
                doc_id=_hash(text),
                text=text,
                source=path,
                doc_type="image_ocr",
            )]
    except ImportError:
        logger.warning("pytesseract or Pillow not installed, cannot OCR images")
    except Exception as e:
        logger.error(f"OCR failed for {path}: {e}")
    return []


def load_pdf_ocr(path: str) -> List[LoadedDocument]:
    """OCR for scanned PDFs - extracts images and runs OCR."""
    try:
        import pdfplumber
        docs = []
        with pdfplumber.open(path) as pdf:
            has_text = False
            for page in pdf.pages:
                text = page.extract_text() or ""
                if text.strip():
                    has_text = True
                    docs.append(LoadedDocument(
                        doc_id=f"{_hash(path)}_p{page.page_number}",
                        text=text,
                        source=path,
                        doc_type="pdf_text",
                        metadata={"page": page.page_number},
                    ))

            if not has_text:
                for page_num, page in enumerate(pdf.pages):
                    try:
                        img = page.to_image(resolution=300)
                        import pytesseract
                        text = pytesseract.image_to_string(img.original)
                        if text.strip():
                            docs.append(LoadedDocument(
                                doc_id=f"{_hash(path)}_ocr_{page_num}",
                                text=text,
                                source=path,
                                doc_type="pdf_ocr",
                                metadata={"page": page_num + 1, "method": "ocr"},
                            ))
                    except Exception:
                        pass
        return docs
    except Exception as e:
        logger.error(f"PDF OCR failed for {path}: {e}")
        return []


def load_jsonl(path: str) -> List[LoadedDocument]:
    import json
    docs = []
    for i, line in enumerate(Path(path).read_text(encoding="utf-8").strip().splitlines()):
        if not line.strip():
            continue
        rec = json.loads(line)
        doc_id = rec.get("id", rec.get("doc_id", f"{_hash(path)}_{i}"))
        text = rec.get("text", str(rec))
        docs.append(LoadedDocument(
            doc_id=doc_id,
            text=text,
            source=path,
            doc_type="jsonl",
        ))
    return docs


def load_json(path: str) -> List[LoadedDocument]:
    import json
    data = json.loads(Path(path).read_text(encoding="utf-8"))
    if isinstance(data, dict):
        return [LoadedDocument(doc_id=k, text=v, source=path, doc_type="json") for k, v in data.items()]
    return []


FORMAT_HANDLERS = {
    ".txt": load_txt,
    ".md": load_markdown,
    ".csv": load_csv,
    ".html": load_html,
    ".htm": load_html,
    ".docx": load_docx,
    ".pptx": load_pptx,
    ".pdf": load_pdf_ocr,
    ".jsonl": load_jsonl,
    ".json": load_json,
    ".png": load_image_ocr,
    ".jpg": load_image_ocr,
    ".jpeg": load_image_ocr,
    ".tiff": load_image_ocr,
    ".bmp": load_image_ocr,
}


def load_file(path: str) -> List[LoadedDocument]:
    """Load any supported file format."""
    p = Path(path)
    if p.is_url() if hasattr(p, 'is_url') else path.startswith(('http://', 'https://')):
        return load_url(path)

    ext = p.suffix.lower()
    handler = FORMAT_HANDLERS.get(ext)
    if handler:
        return handler(path)

    logger.warning(f"Unsupported file format: {ext}")
    return []


def load_directory(directory: str, recursive: bool = True) -> List[LoadedDocument]:
    """Load all supported files from a directory."""
    p = Path(directory)
    pattern = "**/*" if recursive else "*"
    docs = []
    for f in sorted(p.glob(pattern)):
        if f.is_file() and f.suffix.lower() in FORMAT_HANDLERS:
            docs.extend(load_file(str(f)))
    return docs


def chunk_documents(docs: List[LoadedDocument], chunk_size: int = 512, overlap: int = 100) -> Dict[str, str]:
    """Chunk loaded documents into retriever-compatible format."""
    try:
        from langchain_text_splitters import RecursiveCharacterTextSplitter
    except ImportError:
        from langchain.text_splitter import RecursiveCharacterTextSplitter

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=overlap,
        separators=["\n\n", "\n", ". ", "! ", "? ", " ", ""],
    )

    result = {}
    for doc in docs:
        chunks = splitter.split_text(doc.text)
        for i, chunk in enumerate(chunks):
            if len(chunk.strip()) < 50:
                continue
            chunk_id = f"{doc.doc_id}_c{i}"
            result[chunk_id] = chunk
    return result
